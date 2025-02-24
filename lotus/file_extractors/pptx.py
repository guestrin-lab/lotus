"""Slides parser."""

from io import BytesIO
from pathlib import Path

from fsspec import AbstractFileSystem
from llama_index.core.readers.base import BaseReader
from llama_index.core.schema import Document
from llama_index.core.utils import infer_torch_device


class PptxReader(BaseReader):
    """Powerpoint parser.

    Extract text and return it page by page.

    Args:
        should_caption_images (bool): Whether to caption images in the slides.
        caption_model (str): The model to use for image captioning.
        **gen_kwargs: Keyword arguments to pass to the model for image captioning.
    """

    def __init__(
        self, should_caption_images=False, caption_model="nlpconnect/vit-gpt2-image-captioning", **gen_kwargs
    ) -> None:
        try:
            from pptx import Presentation  # noqa
        except ImportError:
            raise ImportError(
                "Please install extra dependencies that are required for "
                "the PptxReader: "
                "`pip install python-pptx`"
            )

        if should_caption_images:
            self._init_caption_images(caption_model)
            self.gen_kwargs = gen_kwargs or {"max_length": 16, "num_beams": 4}
        self.should_caption_images = should_caption_images

    def _init_caption_images(self, caption_model):
        try:
            import torch  # noqa
            from PIL import Image  # noqa
            from transformers import (
                AutoTokenizer,
                VisionEncoderDecoderModel,
                ViTFeatureExtractor,
            )
        except ImportError:
            raise ImportError(
                "Please install extra dependencies that are required for "
                "the PptxReader with Image captions: "
                "`pip install torch transformers python-pptx Pillow`"
            )
        model = VisionEncoderDecoderModel.from_pretrained(caption_model)
        feature_extractor = ViTFeatureExtractor.from_pretrained(caption_model)
        tokenizer = AutoTokenizer.from_pretrained(caption_model)

        self.parser_config = {
            "feature_extractor": feature_extractor,
            "model": model,
            "tokenizer": tokenizer,
        }

    def caption_image(self, image_bytes: bytes) -> str:
        """Generate text caption of image."""
        from PIL import Image

        model = self.parser_config["model"]
        feature_extractor = self.parser_config["feature_extractor"]
        tokenizer = self.parser_config["tokenizer"]

        device = infer_torch_device()
        model.to(device)

        i_image: Image.ImageFile.ImageFile | Image.Image = Image.open(BytesIO(image_bytes))
        if i_image.mode != "RGB":
            i_image = i_image.convert(mode="RGB")

        pixel_values = feature_extractor(images=[i_image], return_tensors="pt").pixel_values
        pixel_values = pixel_values.to(device)

        output_ids = model.generate(pixel_values, **self.gen_kwargs)

        preds = tokenizer.batch_decode(output_ids, skip_special_tokens=True)
        return preds[0].strip()

    def load_data(
        self,
        file: Path,
        extra_info: dict | None = None,
        fs: AbstractFileSystem | None = None,
    ) -> list[Document]:
        """Parse file."""
        from pptx import Presentation

        if fs:
            with fs.open(file) as f:
                presentation = Presentation(f)
        else:
            presentation = Presentation(file)

        docs = []
        for i, slide in enumerate(presentation.slides):
            text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += f"{shape.text}\n"
                if hasattr(shape, "image") and self.should_caption_images:
                    text += f"Image: {self.caption_image(shape.image.blob)}\n\n"
            metadata = {"page_label": i + 1, "file_name": file.name}
            if extra_info is not None:
                metadata.update(extra_info)
            docs.append(Document(text=text, metadata=metadata))

        return docs
